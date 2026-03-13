import os
import sys
import json
import re
import hashlib
import unicodedata
import time
import functools
import logging
import datetime
import urllib.request
import urllib.error
import argparse
import urllib.parse
from dataclasses import dataclass, field

from dotenv import load_dotenv
from openai import AzureOpenAI
from azure.identity import AzureCliCredential, DefaultAzureCredential, get_bearer_token_provider
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import ClientAuthenticationError, HttpResponseError, ResourceNotFoundError

# ---------------------------
# Load environment variables
# ---------------------------
load_dotenv()

# Configure structured logging
logger = logging.getLogger("ingest")
logger.setLevel(logging.INFO)
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter("%(asctime)s %(levelname)s %(message)s")
    handler.setFormatter(formatter)
    logger.addHandler(handler)


AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_EMBEDDING_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT")
AZURE_OPENAI_EMBEDDING_DIMENSIONS = os.getenv("AZURE_OPENAI_EMBEDDING_DIMENSIONS")
AZURE_SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT")
AZURE_SEARCH_INDEX = os.getenv("AZURE_SEARCH_INDEX")
AZURE_SEARCH_KEY = os.getenv("AZURE_SEARCH_KEY")
DOCUMENT_VERSION_DEFAULT = os.getenv("RAG_DOCUMENT_VERSION", "v1.0")
CREATE_NEW_INDEX_VERSION = os.getenv("RAG_CREATE_NEW_INDEX_VERSION", "true").lower() in {"1", "true", "yes"}
DEFAULT_MATTER_ID = os.getenv("RAG_DEFAULT_MATTER_ID", "").strip()

SHAREPOINT_TENANT_ID = os.getenv("SHAREPOINT_TENANT_ID", "").strip()
SHAREPOINT_CLIENT_ID = os.getenv("SHAREPOINT_CLIENT_ID", "").strip()
SHAREPOINT_CLIENT_SECRET = os.getenv("SHAREPOINT_CLIENT_SECRET", "").strip()
SHAREPOINT_DRIVE_ID = os.getenv("SHAREPOINT_DRIVE_ID", "").strip()
SHAREPOINT_FOLDER_PATH = os.getenv("SHAREPOINT_FOLDER_PATH", "").strip()

required_vars = {
    "AZURE_OPENAI_ENDPOINT": AZURE_OPENAI_ENDPOINT,
    "AZURE_OPENAI_EMBEDDING_DEPLOYMENT": AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    "AZURE_SEARCH_ENDPOINT": AZURE_SEARCH_ENDPOINT,
    "AZURE_SEARCH_INDEX": AZURE_SEARCH_INDEX,
}
missing = [name for name, val in required_vars.items() if not val]
if missing:
    print(f"❌ Missing required environment variables: {', '.join(missing)}")
    sys.exit(1)

embedding_dimensions = None
if AZURE_OPENAI_EMBEDDING_DIMENSIONS:
    try:
        embedding_dimensions = int(AZURE_OPENAI_EMBEDDING_DIMENSIONS)
    except ValueError:
        print("❌ AZURE_OPENAI_EMBEDDING_DIMENSIONS must be an integer when provided.")
        sys.exit(1)


def retry(max_attempts: int = 3, initial_delay: float = 1.0, backoff: float = 2.0, exceptions=(Exception,)):
    """Simple retry decorator with exponential backoff."""

    def deco(func):
        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            delay = initial_delay
            last_exc = None
            for attempt in range(1, max_attempts + 1):
                try:
                    return func(*args, **kwargs)
                except exceptions as e:
                    last_exc = e
                    if attempt == max_attempts:
                        logger.error("Giving up after %d attempts: %s", attempt, e)
                        raise
                    logger.warning("Attempt %d failed: %s — retrying in %.1fs", attempt, e, delay)
                    time.sleep(delay)
                    delay *= backoff
            raise last_exc

        return wrapper

    return deco


# A list of (page_number, text) pairs returned by format-specific extractors.
# page_number is 1-indexed. text has already been decoded but not yet cleaned.
PagedText = list[tuple[int, str]]


@dataclass
class SourceDocument:
    source_name: str
    document_id: str
    raw_bytes: bytes
    text: str
    metadata: dict
    paged_text: list = field(default_factory=list)  # list[tuple[int, str]]
    doc_type: str = "text"


# ---------------------------
# Create Azure clients via Entra ID (managed identity / workload identity)
# ---------------------------
use_cli_credential = os.getenv("AZURE_USE_CLI_CREDENTIAL", "").lower() in {"1", "true", "yes"}
credential = AzureCliCredential() if use_cli_credential else DefaultAzureCredential()
openai_token_provider = get_bearer_token_provider(
    credential,
    "https://cognitiveservices.azure.com/.default",
)

openai_client = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    azure_ad_token_provider=openai_token_provider,
    api_version="2024-12-01-preview",
)

search_client = SearchClient(
    endpoint=AZURE_SEARCH_ENDPOINT,
    index_name=AZURE_SEARCH_INDEX,
    credential=AzureKeyCredential(AZURE_SEARCH_KEY) if AZURE_SEARCH_KEY else credential,
)
active_index_name = AZURE_SEARCH_INDEX
checksum_supported = True


def _api_call(method: str, url: str, body: dict | None = None) -> dict:
    if not AZURE_SEARCH_KEY:
        raise RuntimeError("AZURE_SEARCH_KEY is required for index management operations.")

    data = None
    if body is not None:
        data = json.dumps(body).encode("utf-8")

    req = urllib.request.Request(url=url, method=method.upper(), data=data)
    req.add_header("Content-Type", "application/json")
    req.add_header("api-key", AZURE_SEARCH_KEY)
    with urllib.request.urlopen(req, timeout=60) as response:
        payload = response.read().decode("utf-8")
        return json.loads(payload) if payload else {}


def create_versioned_index_if_needed(repo_root: str) -> str:
    global search_client
    global active_index_name

    if not CREATE_NEW_INDEX_VERSION:
        return active_index_name

    if not AZURE_SEARCH_KEY:
        logger.warning("Skipping index version creation because AZURE_SEARCH_KEY is not set.")
        return active_index_name

    timestamp = datetime.datetime.utcnow().strftime("%Y%m%d%H%M%S")
    new_index_name = f"{AZURE_SEARCH_INDEX}-v{timestamp}"
    api_version = "2024-07-01"

    try:
        base_index = _api_call(
            "GET",
            f"{AZURE_SEARCH_ENDPOINT}/indexes/{AZURE_SEARCH_INDEX}?api-version={api_version}",
        )
        base_index["name"] = new_index_name
        _api_call(
            "PUT",
            f"{AZURE_SEARCH_ENDPOINT}/indexes/{new_index_name}?api-version={api_version}",
            base_index,
        )
        logger.info("✅ Created new versioned index: %s", new_index_name)
    except Exception as exc:
        logger.warning("Could not create versioned index; using base index %s. Error: %s", AZURE_SEARCH_INDEX, exc)
        return active_index_name

    active_index_name = new_index_name
    search_client = SearchClient(
        endpoint=AZURE_SEARCH_ENDPOINT,
        index_name=active_index_name,
        credential=AzureKeyCredential(AZURE_SEARCH_KEY) if AZURE_SEARCH_KEY else credential,
    )

    # Update local active pointer shared with the API process.
    data_dir = os.path.join(repo_root, "src", "data")
    os.makedirs(data_dir, exist_ok=True)
    state_path = os.path.join(data_dir, "index-version-state.json")
    known = [active_index_name]
    if os.path.exists(state_path):
        try:
            with open(state_path, "r", encoding="utf-8") as existing:
                parsed = json.load(existing)
            known = parsed.get("KnownIndexes") or parsed.get("knownIndexes") or []
            if not isinstance(known, list):
                known = [active_index_name]
        except Exception:
            known = [active_index_name]

    if active_index_name not in known:
        known.append(active_index_name)

    with open(state_path, "w", encoding="utf-8") as f:
        json.dump(
            {
                "activeIndex": active_index_name,
                "knownIndexes": sorted(set(known)),
                "updatedAt": datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z",
            },
            f,
            indent=2,
        )
    return active_index_name


def compute_checksum(raw_bytes: bytes) -> str:
    digest = hashlib.sha256(raw_bytes).hexdigest()
    return f"sha256-{digest}"


def is_duplicate_checksum_in_index(checksum: str) -> bool:
    escaped = checksum.replace("'", "''")
    results = search_client.search(
        search_text="*",
        filter=f"checksum eq '{escaped}'",
        top=1,
        select=["id"],
    )
    for _ in results:
        return True
    return False


def append_lineage_record(repo_root: str, record: dict) -> None:
    data_dir = os.path.join(repo_root, "src", "data")
    os.makedirs(data_dir, exist_ok=True)
    lineage_path = os.path.join(data_dir, "lineage-records.jsonl")
    with open(lineage_path, "a", encoding="utf-8") as f:
        f.write(json.dumps(record) + "\n")


def validate_openai_access() -> None:
    try:
        embedding_kwargs = {}
        if embedding_dimensions is not None:
            embedding_kwargs["dimensions"] = embedding_dimensions
        openai_client.embeddings.create(
            model=AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
            input="ping",
            **embedding_kwargs,
        )
        logger.info("✅ Azure OpenAI access validated via Entra ID.")
    except ClientAuthenticationError as e:
        logger.error("❌ Azure OpenAI authentication failed.")
        logger.error("Ensure this identity has role 'Cognitive Services OpenAI User' on the OpenAI resource.")
        logger.error("Original error: %s", e)
        sys.exit(1)
    except HttpResponseError as e:
        logger.error("❌ Azure OpenAI request failed.")
        logger.error("Check endpoint/deployment name and role assignments.")
        logger.error("Original error: %s", e)
        sys.exit(1)


def validate_search_access() -> None:
    try:
        count = search_client.get_document_count()
        logger.info("✅ Azure AI Search access validated for index '%s' (documents=%d).", active_index_name, count)
    except ResourceNotFoundError as e:
        logger.error("❌ Search index '%s' not found.", AZURE_SEARCH_INDEX)
        logger.error("Create the index first (IaC or indexer setup) before ingestion.")
        logger.error("Original error: %s", e)
        sys.exit(1)
    except ClientAuthenticationError as e:
        logger.error("❌ Azure AI Search authentication failed.")
        logger.error("Ensure this identity has role 'Search Index Data Contributor' on the Search service.")
        logger.error("Original error: %s", e)
        sys.exit(1)
    except HttpResponseError as e:
        logger.error("❌ Azure AI Search request failed.")
        logger.error("Check endpoint/index name and role assignments.")
        logger.error("Original error: %s", e)
        sys.exit(1)


def clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("utf-8", "ignore").decode("utf-8", "ignore")
    return text.strip()


def chunk_text(text: str, max_chars: int = 2000, overlap: int = 200) -> list[str]:
    if max_chars <= 0:
        return [text]

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if not paragraphs:
        return [text[:max_chars]]

    chunks = []
    current = ""
    for para in paragraphs:
        candidate = para if not current else f"{current}\n\n{para}"
        if len(candidate) <= max_chars:
            current = candidate
            continue

        if current:
            chunks.append(current)
            current = ""

        if len(para) <= max_chars:
            current = para
            continue

        start = 0
        while start < len(para):
            end = start + max_chars
            chunks.append(para[start:end])
            if end >= len(para):
                break
            start = end - overlap if (end - overlap) > start else end

    if current:
        chunks.append(current)

    return chunks


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Production document ingestion for Azure AI Search.")
    parser.add_argument(
        "--source",
        choices=["folder", "sharepoint"],
        default="folder",
        help="Document source type.",
    )
    parser.add_argument(
        "--documents-path",
        default="documents",
        help="Local folder path used when --source folder.",
    )
    parser.add_argument(
        "--max-chars",
        type=int,
        default=2000,
        help="Chunk size in characters.",
    )
    parser.add_argument(
        "--overlap",
        type=int,
        default=200,
        help="Chunk overlap in characters for long segments.",
    )
    parser.add_argument(
        "--batch-size",
        type=int,
        default=100,
        help="Search upload batch size.",
    )
    return parser.parse_args()


def _read_text_from_bytes(raw: bytes, source_name: str) -> str:
    try:
        return clean_text(raw.decode("utf-8"))
    except UnicodeDecodeError:
        logger.warning("Skipping %s because content is not valid UTF-8 text.", source_name)
        return ""


def extract_text_from_pdf(raw: bytes, source_name: str) -> PagedText:
    try:
        import io
        import pdfplumber
    except ImportError:
        logger.error("pdfplumber not installed. Run: pip install pdfplumber")
        return []
    try:
        pages: PagedText = []
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages.append((page_num, text))
        if not any(t.strip() for _, t in pages):
            logger.warning(
                "Skipping %s: no text extracted from any PDF page. "
                "This document may be scanned. Pre-process with OCR before ingesting.",
                source_name,
            )
            return []
        return pages
    except Exception as exc:
        logger.warning("Skipping %s due to PDF extraction error: %s", source_name, exc)
        return []


def extract_text_from_docx(raw: bytes, source_name: str) -> PagedText:
    try:
        import io
        from docx import Document
    except ImportError:
        logger.error("python-docx not installed. Run: pip install python-docx")
        return []
    try:
        doc = Document(io.BytesIO(raw))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or "").lower()
            if style_name.startswith("heading"):
                parts.append(f"\n\n{text}")
            else:
                parts.append(text)
        full_text = "\n".join(parts).strip()
        if not full_text:
            logger.warning("Skipping %s: DOCX contains no extractable text.", source_name)
            return []
        return [(1, full_text)]
    except Exception as exc:
        logger.warning("Skipping %s due to DOCX extraction error: %s", source_name, exc)
        return []


def extract_text_from_xlsx(raw: bytes, source_name: str) -> PagedText:
    try:
        import io
        import openpyxl
    except ImportError:
        logger.error("openpyxl not installed. Run: pip install openpyxl")
        return []
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        sheet_blocks = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    rows_text.append("\t".join(cells))
            if rows_text:
                sheet_blocks.append(f"Sheet: {sheet_name}\n" + "\n".join(rows_text))
        if not sheet_blocks:
            logger.warning("Skipping %s: XLSX contains no data.", source_name)
            return []
        return [(1, "\n\n".join(sheet_blocks))]
    except Exception as exc:
        logger.warning("Skipping %s due to XLSX extraction error: %s", source_name, exc)
        return []


def extract_text_from_pptx(raw: bytes, source_name: str) -> PagedText:
    try:
        import io
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return []
    try:
        prs = Presentation(io.BytesIO(raw))
        pages: PagedText = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            slide_text = "\n".join(texts).strip()
            if slide_text:
                pages.append((slide_num, slide_text))
        if not pages:
            logger.warning("Skipping %s: PPTX contains no extractable text.", source_name)
            return []
        return pages
    except Exception as exc:
        logger.warning("Skipping %s due to PPTX extraction error: %s", source_name, exc)
        return []


def extract_paged_text(raw: bytes, source_name: str) -> tuple[PagedText, str]:
    """Route raw bytes to the correct extractor. Returns (paged_text, doc_type)."""
    lowered = source_name.lower()
    if lowered.endswith(".pdf"):
        return extract_text_from_pdf(raw, source_name), "pdf"
    elif lowered.endswith(".docx"):
        return extract_text_from_docx(raw, source_name), "docx"
    elif lowered.endswith(".xlsx"):
        return extract_text_from_xlsx(raw, source_name), "xlsx"
    elif lowered.endswith(".pptx"):
        return extract_text_from_pptx(raw, source_name), "pptx"
    else:
        text = _read_text_from_bytes(raw, source_name)
        return ([(1, text)], "text") if text else ([], "text")


def _is_supported_text_name(name: str) -> bool:
    return name.lower().endswith((".txt", ".md", ".csv", ".json", ".pdf", ".docx", ".xlsx", ".pptx"))


def load_local_documents(documents_folder: str) -> tuple[list[SourceDocument], dict]:
    logger.info("📂 Local documents folder: %s", documents_folder)
    if not os.path.exists(documents_folder):
        logger.error("Documents folder not found: %s", documents_folder)
        return [], {}

    metadata_map = load_document_metadata(documents_folder)
    docs: list[SourceDocument] = []
    for filename in os.listdir(documents_folder):
        if not _is_supported_text_name(filename):
            logger.warning("Skipping unsupported format: %s", filename)
            continue

        filepath = os.path.join(documents_folder, filename)
        try:
            with open(filepath, "rb") as f:
                raw = f.read()
        except Exception as exc:
            logger.warning("Skipping %s because file could not be read: %s", filename, exc)
            continue

        if not raw:
            logger.warning("Skipping %s because file is empty.", filename)
            continue

        paged, doc_type = extract_paged_text(raw, filename)
        if not paged:
            continue
        text = clean_text("\n\n".join(t for _, t in paged))

        docs.append(
            SourceDocument(
                source_name=filename,
                document_id=filename,
                raw_bytes=raw,
                text=text,
                metadata=metadata_map.get(filename, {}),
                paged_text=paged,
                doc_type=doc_type,
            )
        )

    logger.info("Loaded %d local document(s) for ingestion.", len(docs))
    return docs, metadata_map


@retry(max_attempts=4, initial_delay=1.0, exceptions=(urllib.error.URLError, urllib.error.HTTPError))
def _graph_token() -> str:
    required = {
        "SHAREPOINT_TENANT_ID": SHAREPOINT_TENANT_ID,
        "SHAREPOINT_CLIENT_ID": SHAREPOINT_CLIENT_ID,
        "SHAREPOINT_CLIENT_SECRET": SHAREPOINT_CLIENT_SECRET,
    }
    missing = [k for k, v in required.items() if not v]
    if missing:
        raise RuntimeError(f"Missing required SharePoint auth environment variables: {', '.join(missing)}")

    token_url = f"https://login.microsoftonline.com/{SHAREPOINT_TENANT_ID}/oauth2/v2.0/token"
    form_data = urllib.parse.urlencode(
        {
            "client_id": SHAREPOINT_CLIENT_ID,
            "client_secret": SHAREPOINT_CLIENT_SECRET,
            "scope": "https://graph.microsoft.com/.default",
            "grant_type": "client_credentials",
        }
    ).encode("utf-8")

    req = urllib.request.Request(token_url, method="POST", data=form_data)
    req.add_header("Content-Type", "application/x-www-form-urlencoded")
    with urllib.request.urlopen(req, timeout=60) as response:
        payload = json.loads(response.read().decode("utf-8"))
        token = payload.get("access_token")
        if not token:
            raise RuntimeError("SharePoint token response did not include access_token.")
        return token


@retry(max_attempts=4, initial_delay=1.0, exceptions=(urllib.error.URLError, urllib.error.HTTPError))
def _graph_get(url: str, token: str) -> dict:
    req = urllib.request.Request(url=url, method="GET")
    req.add_header("Authorization", f"Bearer {token}")
    req.add_header("Accept", "application/json")
    with urllib.request.urlopen(req, timeout=60) as response:
        body = response.read().decode("utf-8")
        return json.loads(body) if body else {}


@retry(max_attempts=4, initial_delay=1.0, exceptions=(urllib.error.URLError, urllib.error.HTTPError))
def _download_sharepoint_file(download_url: str) -> bytes:
    req = urllib.request.Request(url=download_url, method="GET")
    with urllib.request.urlopen(req, timeout=120) as response:
        return response.read()


def _sharepoint_children_url(drive_id: str, folder_path: str) -> str:
    encoded_path = urllib.parse.quote(folder_path.strip("/"))
    if encoded_path:
        return f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_path}:/children?$top=999"
    return f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children?$top=999"


def _list_sharepoint_items_recursive(drive_id: str, folder_path: str, token: str) -> list[dict]:
    items: list[dict] = []
    queue = [folder_path]
    while queue:
        current_folder = queue.pop(0)
        next_url = _sharepoint_children_url(drive_id, current_folder)
        while next_url:
            response = _graph_get(next_url, token)
            for item in response.get("value", []):
                if "folder" in item:
                    child_path = item.get("parentReference", {}).get("path", "")
                    child_name = item.get("name", "")
                    relative = child_path.split("root:", 1)[-1].strip("/")
                    next_path = f"{relative}/{child_name}".strip("/")
                    queue.append(next_path)
                    continue
                items.append(item)
            next_url = response.get("@odata.nextLink")
    return items


def load_sharepoint_documents() -> tuple[list[SourceDocument], dict]:
    required = {
        "SHAREPOINT_DRIVE_ID": SHAREPOINT_DRIVE_ID,
    }
    missing = [k for k, v in required.items() if not v]
    if missing:
        logger.error("Missing required SharePoint source variables: %s", ", ".join(missing))
        return [], {}

    token = _graph_token()
    folder = SHAREPOINT_FOLDER_PATH or ""
    logger.info("📂 Reading SharePoint documents from drive %s path '%s'", SHAREPOINT_DRIVE_ID, folder or "/")
    items = _list_sharepoint_items_recursive(SHAREPOINT_DRIVE_ID, folder, token)
    logger.info("Found %d SharePoint file item(s).", len(items))

    docs: list[SourceDocument] = []
    for item in items:
        name = item.get("name", "")
        if not _is_supported_text_name(name):
            logger.warning("Skipping unsupported SharePoint file type: %s", name)
            continue

        download_url = item.get("@microsoft.graph.downloadUrl")
        if not download_url:
            logger.warning("Skipping %s because download URL is missing.", name)
            continue

        try:
            raw = _download_sharepoint_file(download_url)
        except Exception as exc:
            logger.warning("Skipping SharePoint file %s due to download error: %s", name, exc)
            continue

        if not raw:
            logger.warning("Skipping SharePoint file %s because it is empty.", name)
            continue

        paged, doc_type = extract_paged_text(raw, name)
        if not paged:
            continue
        text = clean_text("\n\n".join(t for _, t in paged))

        item_id = item.get("id", "")
        parent_path = item.get("parentReference", {}).get("path", "")
        source_name = f"{parent_path}/{name}".replace("/drive/root:", "").lstrip("/")
        docs.append(
            SourceDocument(
                source_name=source_name or name,
                document_id=item_id or name,
                raw_bytes=raw,
                text=text,
                metadata={
                    "sourceType": "sharepoint",
                    "sharepointItemId": item_id,
                    "sharepointDriveId": SHAREPOINT_DRIVE_ID,
                    "sharepointPath": source_name,
                    "sharepointLastModified": item.get("lastModifiedDateTime", ""),
                    "sharepointSize": item.get("size"),
                },
                paged_text=paged,
                doc_type=doc_type,
            )
        )

    logger.info("Loaded %d SharePoint text document(s) for ingestion.", len(docs))
    return docs, {}


@retry(max_attempts=3, initial_delay=1.0)
def create_embeddings(input_data):
    embedding_kwargs = {}
    if embedding_dimensions is not None:
        embedding_kwargs["dimensions"] = embedding_dimensions
    return openai_client.embeddings.create(
        model=AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
        input=input_data,
        **embedding_kwargs,
    )


@retry(max_attempts=3, initial_delay=1.0)
def upload_documents_with_retry(documents):
    return search_client.upload_documents(documents)


def load_document_metadata(documents_folder: str) -> dict:
    metadata_path = os.path.join(documents_folder, "metadata.json")
    if not os.path.exists(metadata_path):
        logger.warning("No metadata file found at %s. Create it to enforce matter-level ingestion.", metadata_path)
        return {}

    try:
        with open(metadata_path, "r", encoding="utf-8") as f:
            metadata = json.load(f)
    except Exception as e:
        logger.error("Failed to parse metadata file '%s': %s", metadata_path, e)
        return {}

    if not isinstance(metadata, dict):
        logger.error("Metadata file must be a JSON object keyed by filename.")
        return {}

    logger.info("Loaded metadata for %d document(s).", len(metadata))
    return metadata


def extract_matter_id_from_text(text: str) -> str:
    # Prefer explicit header-style declarations near the top of the document.
    head = text[:8000]
    labeled_patterns = [
        r"(?im)^\s*matter(?:\s+id|\s+number)?\s*[:#-]\s*([A-Za-z0-9][A-Za-z0-9._-]{1,63})\s*$",
        r"(?im)^\s*file\s*matter\s*[:#-]\s*([A-Za-z0-9][A-Za-z0-9._-]{1,63})\s*$",
    ]
    for pattern in labeled_patterns:
        match = re.search(pattern, head)
        if match:
            return match.group(1).strip()

    # Fallback: look for canonical MATTER-* token anywhere in the leading content.
    token_match = re.search(r"(?i)\b(MATTER-[A-Za-z0-9._-]{1,63})\b", head)
    if token_match:
        return token_match.group(1).strip()

    return ""


def extract_matter_id_from_filename(filename: str) -> str:
    match = re.search(r"(?i)\b(MATTER-[A-Za-z0-9._-]{1,63})\b", filename)
    if not match:
        return ""
    return match.group(1).strip()


def resolve_matter_id(filename: str, text: str, doc_metadata: dict) -> tuple[str, str]:
    metadata_matter_id = (doc_metadata.get("matterId") or "").strip() if isinstance(doc_metadata, dict) else ""
    if metadata_matter_id:
        return metadata_matter_id, "metadata"

    text_matter_id = extract_matter_id_from_text(text)
    if text_matter_id:
        return text_matter_id, "document-text"

    filename_matter_id = extract_matter_id_from_filename(filename)
    if filename_matter_id:
        return filename_matter_id, "filename"

    return "", ""


def extract_labeled_value(text: str, labels: list[str], max_chars: int = 8000) -> str:
    head = text[:max_chars]
    for label in labels:
        pattern = rf"(?im)^\s*{label}\s*[:#-]\s*(.+?)\s*$"
        match = re.search(pattern, head)
        if match:
            value = match.group(1).strip()
            if value:
                return value
    return ""


def normalize_confidentiality_level(value: str) -> str:
    if not value:
        return ""
    val = value.strip().lower()
    mapping = {
        "public": "Public",
        "internal": "Internal",
        "confidential": "Confidential",
        "restricted": "Restricted",
        "highly confidential": "Highly Confidential",
    }
    if val in mapping:
        return mapping[val]

    # Keyword-based fallback if the line includes extra words.
    if "highly confidential" in val:
        return "Highly Confidential"
    if "restricted" in val:
        return "Restricted"
    if "confidential" in val:
        return "Confidential"
    if "internal" in val:
        return "Internal"
    if "public" in val:
        return "Public"
    return value.strip()


def resolve_practice_area(text: str, doc_metadata: dict) -> tuple[str, str]:
    metadata_value = (doc_metadata.get("practiceArea") or "").strip() if isinstance(doc_metadata, dict) else ""
    if metadata_value:
        return metadata_value, "metadata"

    text_value = extract_labeled_value(text, ["Practice\\s*Area", "Area\\s*of\\s*Law"])
    if text_value:
        return text_value, "document-text"

    return "", ""


def resolve_client(text: str, doc_metadata: dict) -> tuple[str, str]:
    metadata_value = (doc_metadata.get("client") or "").strip() if isinstance(doc_metadata, dict) else ""
    if metadata_value:
        return metadata_value, "metadata"

    text_value = extract_labeled_value(text, ["Client", "Customer", "Represented\\s*Party"])
    if text_value:
        return text_value, "document-text"

    return "", ""


def resolve_confidentiality_level(text: str, doc_metadata: dict) -> tuple[str, str]:
    metadata_value = (doc_metadata.get("confidentialityLevel") or "").strip() if isinstance(doc_metadata, dict) else ""
    if metadata_value:
        return normalize_confidentiality_level(metadata_value), "metadata"

    text_value = extract_labeled_value(text, ["Confidentiality\\s*Level", "Confidentiality", "Classification"])
    if text_value:
        return normalize_confidentiality_level(text_value), "document-text"

    return "", ""


def resolve_acl(doc_metadata: dict) -> tuple[list[str], list[str]]:
    """
    Resolve document ACL (access-control lists) from per-document metadata.

    Returns:
        (allowedUsers, allowedGroups)

    allowedUsers:  list of UPNs or email addresses permitted to read this document
                   (matched against the preferred_username / email claim in Entra ID tokens).
    allowedGroups: list of Entra ID group object IDs (or display names) permitted to read
                   this document (matched against the groups claim in Entra ID tokens).

    Default behaviour when no ACL is specified: grant read access to the "all-lawyers"
    group so that newly ingested documents are accessible to all firm staff until explicit
    ACLs are applied.
    """
    if not isinstance(doc_metadata, dict):
        return [], ["all-lawyers"]

    allowed_users = doc_metadata.get("allowedUsers", [])
    if not isinstance(allowed_users, list):
        allowed_users = [allowed_users] if allowed_users else []
    allowed_users = [u.strip() for u in allowed_users if u and str(u).strip()]

    allowed_groups = doc_metadata.get("allowedGroups", [])
    if not isinstance(allowed_groups, list):
        allowed_groups = [allowed_groups] if allowed_groups else []
    allowed_groups = [g.strip() for g in allowed_groups if g and str(g).strip()]

    # Default: broad firm-wide read access when no explicit ACL has been assigned.
    if not allowed_users and not allowed_groups:
        allowed_groups = ["all-lawyers"]

    return allowed_users, allowed_groups


def ensure_acl_fields_in_index() -> None:
    """
    Add allowedUsers and allowedGroups Collection(Edm.String) fields to the active
    index if they are not already present.  These fields are required for document-level
    ACL security trimming in the API layer.

    The fields are marked filterable=True and retrievable=False.  They are stored in
    the index for filter evaluation only and are never returned in query results, which
    prevents ACL metadata from leaking through the API response.

    Requires AZURE_SEARCH_KEY to be set (used for index management REST calls).
    """
    if not AZURE_SEARCH_KEY:
        logger.warning(
            "AZURE_SEARCH_KEY not set; cannot auto-update index schema. "
            "Manually add 'allowedUsers' and 'allowedGroups' Collection(Edm.String) "
            "filterable fields to index '%s' to enable document-level ACL filtering.",
            active_index_name,
        )
        return

    api_version = "2024-07-01"
    try:
        schema = _api_call(
            "GET",
            f"{AZURE_SEARCH_ENDPOINT}/indexes/{active_index_name}?api-version={api_version}",
        )
    except Exception as exc:
        logger.warning("Could not retrieve index schema for ACL field check: %s", exc)
        return

    existing_names = {f["name"] for f in schema.get("fields", [])}

    new_fields = []
    for field_name in ("allowedUsers", "allowedGroups"):
        if field_name not in existing_names:
            new_fields.append(
                {
                    "name": field_name,
                    "type": "Collection(Edm.String)",
                    "filterable": True,
                    "retrievable": False,
                    "searchable": False,
                    "sortable": False,
                    "facetable": False,
                }
            )

    if not new_fields:
        logger.info("ACL fields (allowedUsers, allowedGroups) already present in index '%s'.", active_index_name)
        return

    schema["fields"].extend(new_fields)
    try:
        _api_call(
            "PUT",
            f"{AZURE_SEARCH_ENDPOINT}/indexes/{active_index_name}?api-version={api_version}",
            schema,
        )
        added = [f["name"] for f in new_fields]
        logger.info("Added ACL fields %s to index '%s'.", added, active_index_name)
    except Exception as exc:
        logger.warning(
            "Could not update index schema with ACL fields: %s. "
            "Manually add 'allowedUsers' and 'allowedGroups' Collection(Edm.String) "
            "filterable fields before ingesting.",
            exc,
        )


def ingest_documents(source: str, documents_path: str, max_chars: int, overlap: int, batch_size: int) -> None:
    base_dir = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(base_dir)
    run_ingestion_timestamp = datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
    global checksum_supported
    if source == "folder":
        full_documents_path = documents_path if os.path.isabs(documents_path) else os.path.join(repo_root, documents_path)
        source_docs, _ = load_local_documents(full_documents_path)
    else:
        source_docs, _ = load_sharepoint_documents()

    seen_checksums = set()

    documents_to_upload = []
    lineage_buffer = []

    for source_doc in source_docs:
        filename = source_doc.source_name
        doc_metadata = source_doc.metadata
        raw = source_doc.raw_bytes
        text = source_doc.text
        logger.info("🔍 Processing document: %s", filename)

        checksum = compute_checksum(raw)
        if checksum in seen_checksums:
            logger.info("Skipping %s because duplicate checksum detected in this run (%s).", filename, checksum)
            continue
        seen_checksums.add(checksum)

        duplicate_in_index = False
        if checksum_supported:
            try:
                duplicate_in_index = is_duplicate_checksum_in_index(checksum)
            except Exception as exc:
                message = str(exc).lower()
                if "checksum" in message or "checksum" in repr(exc).lower():
                    checksum_supported = False
                    logger.warning("Checksum field missing in index; skipping further duplicate checks.")
                else:
                    logger.warning("Checksum duplicate check failed for %s: %s", filename, exc)
        if duplicate_in_index:
            logger.info("Skipping %s because duplicate checksum already exists in index (%s).", filename, checksum)
            continue

        logger.info("🧪 Text length: %d", len(text))
        if not text:
            logger.warning("Empty after cleaning. Skipping %s", filename)
            continue

        matter_id, matter_id_source = resolve_matter_id(filename, text, doc_metadata)
        if not matter_id and DEFAULT_MATTER_ID:
            matter_id = DEFAULT_MATTER_ID
            matter_id_source = "default-env"
        if not matter_id:
            logger.warning("Skipping %s because matterId was not found in metadata, document text, or filename.", filename)
            continue
        logger.info("🧾 matterId resolved for %s: %s (%s)", filename, matter_id, matter_id_source)
        practice_area, practice_area_source = resolve_practice_area(text, doc_metadata)
        client, client_source = resolve_client(text, doc_metadata)
        confidentiality_level, confidentiality_level_source = resolve_confidentiality_level(text, doc_metadata)
        allowed_users, allowed_groups = resolve_acl(doc_metadata)
        if practice_area:
            logger.info("🧾 practiceArea resolved for %s: %s (%s)", filename, practice_area, practice_area_source)
        if client:
            logger.info("🧾 client resolved for %s: %s (%s)", filename, client, client_source)
        if confidentiality_level:
            logger.info(
                "🧾 confidentialityLevel resolved for %s: %s (%s)",
                filename,
                confidentiality_level,
                confidentiality_level_source,
            )
        logger.info("🔐 ACL resolved for %s: users=%s groups=%s", filename, allowed_users, allowed_groups)

        all_chunks: list[tuple[int, str]] = []
        for page_num, page_text in source_doc.paged_text:
            page_clean = clean_text(page_text)
            if not page_clean:
                continue
            for c in chunk_text(page_clean, max_chars=max_chars, overlap=overlap):
                all_chunks.append((page_num, c))

        if not all_chunks:
            logger.warning("Empty after cleaning. Skipping %s", filename)
            continue

        logger.info(
            "🧠 Creating embeddings for %d chunk(s) across %d page(s)...",
            len(all_chunks),
            len(source_doc.paged_text),
        )

        chunk_strings = [c for _, c in all_chunks]
        try:
            response = create_embeddings(chunk_strings)
            embeddings = [item.embedding for item in response.data]
        except Exception as e:
            logger.error("Failed to create embeddings for %s: %s", filename, e)
            continue

        document_id = (doc_metadata.get("documentId") or source_doc.document_id or filename).strip() if isinstance(doc_metadata, dict) else filename
        document_version = (doc_metadata.get("documentVersion") or DOCUMENT_VERSION_DEFAULT).strip() if isinstance(doc_metadata, dict) else DOCUMENT_VERSION_DEFAULT

        for idx, ((page_num, chunk_text_val), emb) in enumerate(zip(all_chunks, embeddings)):
            chunk_id = f"{document_id}::chunk-{idx}"
            documents_to_upload.append(
                {
                    "id": chunk_id,
                    "content": chunk_text_val,
                    "contentVector": emb,
                    "source": filename,
                    "sourceFile": filename,
                    "page": page_num,
                    "chunk_index": idx,
                    "documentType": source_doc.doc_type,
                    "matterId": matter_id,
                    "practiceArea": practice_area,
                    "client": client,
                    "confidentialityLevel": confidentiality_level,
                    "documentVersion": document_version,
                    "ingestionTimestamp": run_ingestion_timestamp,
                    "checksum": checksum,
                    "allowedUsers": allowed_users,
                    "allowedGroups": allowed_groups,
                }
            )

        lineage_buffer.append(
            {
                "documentId": document_id,
                "sourceFile": filename,
                "checksum": checksum,
                "ingestionTimestamp": run_ingestion_timestamp,
                "indexedChunks": len(all_chunks),
                "documentType": source_doc.doc_type,
                "indexVersion": active_index_name,
                "eventType": "ingest",
            }
        )

    if not documents_to_upload:
        logger.warning("⚠️ Nothing to upload.")
        return

    logger.info("⬆️ Uploading to Azure Search in batches...")
    total = len(documents_to_upload)
    for i in range(0, total, batch_size):
        batch = documents_to_upload[i : i + batch_size]
        try:
            upload_documents_with_retry(batch)
            logger.info("✅ Uploaded batch %d (%d documents)", i // batch_size + 1, len(batch))
        except Exception as e:
            logger.error("❌ Upload failed for batch starting at %d: %s", i, e)

    for record in lineage_buffer:
        append_lineage_record(repo_root, record)


if __name__ == "__main__":
    args = parse_args()
    repo_root_for_indexing = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    create_versioned_index_if_needed(repo_root_for_indexing)
    validate_openai_access()
    validate_search_access()
    ensure_acl_fields_in_index()
    ingest_documents(
        source=args.source,
        documents_path=args.documents_path,
        max_chars=args.max_chars,
        overlap=args.overlap,
        batch_size=args.batch_size,
    )
